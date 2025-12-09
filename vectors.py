# vectors.py - Document processing and embeddings management with enhanced security and performance

import os
import logging
import hashlib
import time
import uuid
import threading
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
from functools import wraps
from langchain_community.document_loaders import PyPDFLoader, UnstructuredPDFLoader, TextLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import Qdrant
from langchain_core.documents import Document
from langchain_core.documents.base import Document

# Vector database imports
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams, PointStruct

# Additional document loaders
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logging.warning("python-docx not available. Install with: pip install python-docx")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available. Install with: pip install python-pptx")

from dotenv import load_dotenv

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)
load_dotenv()

# FIXED: Realistic Cost and Security limits for production use
MAX_DOCUMENT_SIZE = 50 * 1024 * 1024  # 50MB
MAX_CHUNKS_PER_DOCUMENT = 1000  # Reduced for better performance
MAX_PROCESSING_TIME = 300  # 5 minutes
ALLOWED_FILE_TYPES = {'.pdf', '.txt', '.docx', '.pptx'}
MAX_CONCURRENT_PROCESSING = 2  # Reduced to prevent resource exhaustion

# Thread-safe processing lock
processing_lock = threading.Lock()
active_sessions = {}

def retry_on_failure(max_retries=3, delay=1):
    """Decorator for retrying failed operations"""
    def decorator(func):
        @wraps(func)
        def wrapper(*args, **kwargs):
            for attempt in range(max_retries):
                try:
                    return func(*args, **kwargs)
                except Exception as e:
                    if attempt == max_retries - 1:
                        raise e
                    logger.warning(f"Attempt {attempt + 1} failed: {e}. Retrying in {delay}s...")
                    time.sleep(delay * (attempt + 1))
            return None
        return wrapper
    return decorator

def validate_file_security(file_path: str) -> Tuple[bool, str]:
    """Enhanced security validation for uploaded files"""
    try:
        path = Path(file_path)
        
        # Check file existence
        if not path.exists():
            return False, "File does not exist"
        
        # Check file size
        file_size = path.stat().st_size
        if file_size > MAX_DOCUMENT_SIZE:
            return False, f"File too large: {file_size / (1024*1024):.1f}MB > {MAX_DOCUMENT_SIZE / (1024*1024)}MB"
        
        # Check file extension
        if path.suffix.lower() not in ALLOWED_FILE_TYPES:
            return False, f"File type not allowed: {path.suffix}"
        
        # Check for suspicious file patterns
        suspicious_patterns = ['../', '..\\', '<script', 'javascript:', 'data:', 'vbscript:', 'on load=', 'on click=']
        if any(pattern in path.name.lower() for pattern in suspicious_patterns):
            return False, "Suspicious file name detected"
        
        # Check for very small files (likely empty or corrupted)
        if file_size < 50:  # 50 bytes minimum
            return False, "File too small (likely empty or corrupted)"
        
        # Check file extension consistency (basic anti-spoofing)
        try:
            with open(file_path, 'rb') as f:
                file_header = f.read(8)
                
            # Basic file signature checks
            pdf_signatures = [b'%PDF-', b'%FDF-']
            docx_signatures = [b'PK\x03\x04', b'PK\x05\x06', b'PK\x07\x08']
            
            if path.suffix.lower() == '.pdf':
                if not any(file_header.startswith(sig) for sig in pdf_signatures):
                    return False, "File content doesn't match PDF format"
            elif path.suffix.lower() in ['.docx', '.pptx']:
                if not any(file_header.startswith(sig) for sig in docx_signatures):
                    return False, f"File content doesn't match {path.suffix.upper()} format"
        except:
            # If we can't read the file, it's probably corrupted
            return False, "Unable to read file - may be corrupted"
        
        return True, "File validation passed"
        
    except Exception as e:
        return False, f"Validation error: {str(e)}"

class EmbeddingsManager:
    """FIXED: Enhanced embeddings manager with persistent collection and better error handling"""
    
    # FIXED: Use a single, persistent collection name for all documents
    DEFAULT_COLLECTION_NAME = "Legal_documents"
    
    def __init__(
        self,
        model_name: str = "BAAI/bge-small-en",
        device: str = "cpu",
        encode_kwargs: dict = None,
        qdrant_url: str = None,
        collection_name: str = None,
        chunk_size: int = 1200,
        chunk_overlap: int = 300,
        max_chunks: int = MAX_CHUNKS_PER_DOCUMENT
    ):
        """Initialize with improved validation and persistent collection"""
        
        # Generate unique session ID for this instance
        self.session_id = str(uuid.uuid4())[:8]
        
        # Validate and sanitize inputs
        self.model_name = self._sanitize_model_name(model_name)
        self.device = device if device in ["cpu", "cuda"] else "cpu"
        self.encode_kwargs = encode_kwargs or {"normalize_embeddings": True}
        self.qdrant_url = qdrant_url or os.getenv('QDRANT_URL')
        
        # FIXED: Use persistent collection name - all sessions share the same collection
        # This ensures documents persist across sessions and are not lost
        if collection_name:
            self.collection_name = self._sanitize_collection_name(collection_name)
        else:
            # Use the default persistent collection name
            self.collection_name = self.DEFAULT_COLLECTION_NAME
            
        self.api_key = os.getenv('QDRANT_API_KEY')
        
        # Validate and set processing parameters
        self.chunk_size = max(500, min(chunk_size, 2000))
        self.chunk_overlap = max(50, min(chunk_overlap, self.chunk_size // 2))
        self.max_chunks = max_chunks
        
        # Processing statistics with thread safety
        self.stats_lock = threading.Lock()
        self.stats = {
            'documents_processed': 0,
            'total_chunks_created': 0,
            'processing_time': 0,
            'errors_encountered': 0,
            'session_id': self.session_id,
            'files_processed': []  # Track processed file hashes
        }
        
        # Track this session
        with processing_lock:
            active_sessions[self.session_id] = {
                'start_time': time.time(),
                'collection_name': self.collection_name,
                'stats': self.stats
            }
        
        # Initialize embeddings model
        try:
            self.embeddings = HuggingFaceEmbeddings(
                model_name=self.model_name,
                model_kwargs={"device": self.device},
                encode_kwargs=self.encode_kwargs
            )
            logger.info(f"Session {self.session_id}: Initialized embeddings model: {self.model_name}")
        except Exception as e:
            logger.error(f"Session {self.session_id}: Failed to initialize embeddings: {e}")
            raise
        
        # Initialize Qdrant client
        try:
            self.qdrant_client = QdrantClient(
                url=self.qdrant_url,
                api_key=self.api_key,
                prefer_grpc=False
            )
            logger.info(f"Session {self.session_id}: Connected to Qdrant at {self.qdrant_url}")
        except Exception as e:
            logger.error(f"Session {self.session_id}: Failed to connect to Qdrant: {e}")
            raise
        
        # FIXED: Initialize or verify collection without clearing existing data
        # The collection will be created only if it doesn't exist
        # Existing data in the collection will be preserved
        self._initialize_collection()
    
    def _sanitize_model_name(self, model_name: str) -> str:
        """Sanitize model name for security"""
        allowed_models = [
            "BAAI/bge-small-en",
            "BAAI/bge-base-en",
            "sentence-transformers/all-MiniLM-L6-v2",
            "sentence-transformers/all-mpnet-base-v2"
        ]
        return model_name if model_name in allowed_models else "BAAI/bge-small-en"
    
    def _sanitize_collection_name(self, collection_name: str) -> str:
        """Sanitize collection name for security"""
        # Remove special characters, keep alphanumeric and underscores
        sanitized = ''.join(c if c.isalnum() or c == '_' else '_' for c in collection_name)
        # Limit length
        sanitized = sanitized[:64]
        # Ensure not empty
        return sanitized if sanitized else self.DEFAULT_COLLECTION_NAME
    
    def _initialize_collection(self) -> None:
        """
        FIXED: Initialize or verify Qdrant collection
        - Creates collection only if it doesn't exist
        - Does NOT clear or delete existing data
        - Preserves all previously uploaded documents
        """
        try:
            collections = self.qdrant_client.get_collections().collections
            collection_names = [c.name for c in collections]
            
            if self.collection_name not in collection_names:
                # Create new collection only if it doesn't exist
                self.qdrant_client.create_collection(
                    collection_name=self.collection_name,
                    vectors_config=VectorParams(
                        size=384,  # BGE-small embedding dimension
                        distance=Distance.COSINE
                    )
                )
                logger.info(f"Session {self.session_id}: Created new collection: {self.collection_name}")
            else:
                # Collection exists - do NOT clear it, just use it
                # This preserves all existing documents
                logger.info(f"Session {self.session_id}: Using existing collection: {self.collection_name} (existing data preserved)")
                
        except Exception as e:
            logger.error(f"Session {self.session_id}: Collection initialization failed: {e}")
            raise
    
    def _calculate_file_hash(self, file_path: str) -> str:
        """Calculate SHA-256 hash of file for duplicate detection"""
        sha256_hash = hashlib.sha256()
        with open(file_path, "rb") as f:
            for byte_block in iter(lambda: f.read(4096), b""):
                sha256_hash.update(byte_block)
        return sha256_hash.hexdigest()
    
    def _load_document(self, file_path: str) -> List[Document]:
        """Load document using appropriate loader"""
        path = Path(file_path)
        extension = path.suffix.lower()
        
        try:
            if extension == '.pdf':
                # Try PyPDFLoader first, fall back to UnstructuredPDFLoader
                try:
                    loader = PyPDFLoader(file_path)
                    documents = loader.load()
                except Exception as e:
                    logger.warning(f"PyPDFLoader failed: {e}. Trying UnstructuredPDFLoader...")
                    loader = UnstructuredPDFLoader(file_path)
                    documents = loader.load()
                    
            elif extension == '.txt':
                loader = TextLoader(file_path, encoding='utf-8')
                documents = loader.load()
                
            elif extension == '.docx' and DOCX_AVAILABLE:
                # Custom DOCX loader
                documents = self._load_docx(file_path)
                
            elif extension == '.pptx' and PPTX_AVAILABLE:
                # Custom PPTX loader
                documents = self._load_pptx(file_path)
                
            else:
                raise ValueError(f"Unsupported file type: {extension}")
            
            if not documents:
                raise ValueError("No content could be extracted from the document")
            
            return documents
            
        except Exception as e:
            logger.error(f"Document loading failed for {file_path}: {e}")
            raise
    
    def _load_docx(self, file_path: str) -> List[Document]:
        """Load DOCX document"""
        doc = DocxDocument(file_path)
        full_text = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
        
        content = "\n".join(full_text)
        
        return [Document(
            page_content=content,
            metadata={
                "source": file_path,
                "file_type": "docx",
                "file_name": Path(file_path).name
            }
        )]
    
    def _load_pptx(self, file_path: str) -> List[Document]:
        """Load PPTX document"""
        prs = Presentation(file_path)
        full_text = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    full_text.append(shape.text)
        
        content = "\n".join(full_text)
        
        return [Document(
            page_content=content,
            metadata={
                "source": file_path,
                "file_type": "pptx",
                "file_name": Path(file_path).name
            }
        )]
    
    def _split_documents(self, documents: List[Document]) -> List[Document]:
        """Split documents into chunks with enhanced metadata"""
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=self.chunk_overlap,
            separators=["\n\n", "\n", ". ", " ", ""],
            length_function=len
        )
        
        splits = text_splitter.split_documents(documents)
        
        # Limit chunks and add enhanced metadata
        if len(splits) > self.max_chunks:
            logger.warning(f"Document has {len(splits)} chunks, limiting to {self.max_chunks}")
            splits = splits[:self.max_chunks]
        
        # Add enhanced metadata to each chunk
        for i, split in enumerate(splits):
            split.metadata.update({
                'chunk_index': i,
                'total_chunks': len(splits),
                'chunk_size': len(split.page_content),
                'word_count': len(split.page_content.split())
            })
        
        return splits
    
    @retry_on_failure(max_retries=3, delay=2)
    def create_embeddings(self, file_path: str) -> str:
        """
        FIXED: Create and store embeddings for document
        - Adds documents to the persistent collection
        - Does NOT recreate or clear the collection
        - Preserves all existing documents
        """
        start_time = time.time()
        
        try:
            # Security validation
            is_valid, message = validate_file_security(file_path)
            if not is_valid:
                raise SecurityError(f"Security validation failed: {message}")
            
            # Calculate file hash for duplicate detection
            file_hash = self._calculate_file_hash(file_path)
            
            # Check if file already processed (optional - can be removed if you want to allow re-processing)
            with self.stats_lock:
                if file_hash in self.stats['files_processed']:
                    logger.warning(f"Session {self.session_id}: File already processed: {Path(file_path).name}")
                    return f"⚠️ File '{Path(file_path).name}' was already processed in this session"
            
            # Load document
            documents = self._load_document(file_path)
            
            # Add file hash to metadata
            for doc in documents:
                doc.metadata['file_hash'] = file_hash
                doc.metadata['processed_time'] = time.time()
            
            # Split documents
            splits = self._split_documents(documents)
            
            if not splits:
                raise ValueError("No text chunks created from document")
            
            # FIXED: Store embeddings in the persistent collection
            # This adds to existing data, does NOT clear or recreate the collection
            success = False
            methods = [
                ("standard_langchain", self._create_embeddings_method1),
                ("manual_upsert", self._create_embeddings_method2),
                ("batch_texts", self._create_embeddings_method3)
            ]
            
            for method_name, method_func in methods:
                try:
                    method_func(splits)
                    logger.info(f"Session {self.session_id}: Embeddings created using {method_name}")
                    success = True
                    break
                except Exception as e:
                    logger.warning(f"Session {self.session_id}: {method_name} failed: {e}")
                    continue
            
            if not success:
                raise ConnectionError("All embedding creation methods failed")
            
            # Update statistics thread-safely
            processing_time = time.time() - start_time
            with self.stats_lock:
                self.stats.update({
                    'documents_processed': self.stats['documents_processed'] + 1,
                    'total_chunks_created': self.stats['total_chunks_created'] + len(splits),
                    'processing_time': self.stats['processing_time'] + processing_time,
                    'files_processed': self.stats['files_processed'] + [file_hash]
                })
            
            return self._generate_success_message(file_path, splits, processing_time)
            
        except Exception as e:
            with self.stats_lock:
                self.stats['errors_encountered'] += 1
            logger.error(f"Session {self.session_id}: Embedding creation failed: {e}")
            raise Exception(f"Embedding creation failed: {e}")

    def _create_embeddings_method1(self, splits: List[Document]) -> None:
        """Method 1: Standard from_documents"""
        Qdrant.from_documents(
            documents=splits,
            embedding=self.embeddings,
            url=self.qdrant_url,
            api_key=self.api_key,
            collection_name=self.collection_name,
            prefer_grpc=False,
        )

    def _create_embeddings_method2(self, splits: List[Document]) -> None:
        """Method 2: Manual upsert with batching for better performance"""
        texts = [doc.page_content for doc in splits]
        batch_size = 25  # Reduced batch size for stability
        
        for i in range(0, len(texts), batch_size):
            batch_texts = texts[i:i+batch_size]
            batch_docs = splits[i:i+batch_size]
            
            try:
                embeddings_list = self.embeddings.embed_documents(batch_texts)
                
                points = []
                for j, (doc, embedding) in enumerate(zip(batch_docs, embeddings_list)):
                    point_id = f"{doc.metadata.get('file_hash', 'unknown')}_{i+j}_{int(time.time())}"
                    point = PointStruct(
                        id=point_id,
                        vector=embedding,
                        payload={
                            "page_content": doc.page_content,
                            **doc.metadata
                        }
                    )
                    points.append(point)
                
                self.qdrant_client.upsert(
                    collection_name=self.collection_name,
                    points=points
                )
                
                logger.debug(f"Session {self.session_id}: Processed batch {i//batch_size + 1}/{(len(texts)-1)//batch_size + 1}")
                
            except Exception as e:
                logger.error(f"Session {self.session_id}: Batch {i//batch_size + 1} failed: {e}")
                raise

    def _create_embeddings_method3(self, splits: List[Document]) -> None:
        """Method 3: Batch texts"""
        texts = [doc.page_content for doc in splits]
        metadatas = [doc.metadata for doc in splits]
        
        Qdrant.from_texts(
            texts=texts,
            embedding=self.embeddings,
            metadatas=metadatas,
            url=self.qdrant_url,
            api_key=self.api_key,
            collection_name=self.collection_name,
            prefer_grpc=False,
        )

    def _generate_success_message(self, file_path: str, splits: List[Document], processing_time: float) -> str:
        """Generate detailed success message"""
        file_name = Path(file_path).name
        total_chunks = len(splits)
        avg_chunk_size = sum(len(split.page_content) for split in splits) / len(splits)
        total_words = sum(split.metadata.get('word_count', 0) for split in splits)
        
        return (
            f"✅ Successfully processed '{file_name}'!\n\n"
            f"📊 Processing Summary:\n"
            f"• Created {total_chunks} text chunks\n"
            f"• Average chunk size: {avg_chunk_size:.0f} characters\n"
            f"• Total words processed: {total_words:,}\n"
            f"• Processing time: {processing_time:.2f} seconds\n"
            f"• Collection: {self.collection_name}\n"
            f"• Session: {self.session_id}\n"
            f"• Security validation: Passed\n"
            f"• Duplicate check: Verified unique\n"
            f"• Data persistence: Enabled (documents preserved across sessions)"
        )

    def get_stats(self) -> Dict[str, Any]:
        """Get processing statistics thread-safely"""
        with self.stats_lock:
            stats_copy = self.stats.copy()
        
        return {
            **stats_copy,
            "collection_name": self.collection_name,
            "max_chunks_limit": self.max_chunks,
            "chunk_size": self.chunk_size,
            "chunk_overlap": self.chunk_overlap,
            "session_active_time": time.time() - active_sessions.get(self.session_id, {}).get('start_time', time.time()),
            "files_processed_count": len(stats_copy.get('files_processed', []))
        }

    def clear_collection(self) -> None:
        """
        FIXED: Clear all documents from the persistent collection
        WARNING: This will delete ALL documents from all sessions!
        Only use when you want to completely reset the document database.
        """
        try:
            self.qdrant_client.delete_collection(self.collection_name)
            self._initialize_collection()
            
            # Reset file processing tracking
            with self.stats_lock:
                self.stats['files_processed'] = []
            
            logger.warning(f"Session {self.session_id}: Collection CLEARED: {self.collection_name} (all documents deleted)")
        except Exception as e:
            logger.error(f"Session {self.session_id}: Failed to clear collection: {e}")
            raise

    def cleanup_session(self) -> None:
        """
        FIXED: Clean up session resources WITHOUT deleting the collection
        The collection is persistent and should remain available for future sessions
        """
        try:
            # Remove from active sessions
            with processing_lock:
                if self.session_id in active_sessions:
                    del active_sessions[self.session_id]
            
            logger.info(f"Session {self.session_id}: Session cleaned up (collection preserved)")
            
        except Exception as e:
            logger.error(f"Session {self.session_id}: Session cleanup failed: {e}")

    def __del__(self):
        """
        FIXED: Cleanup when object is destroyed - but DON'T delete the collection
        The collection is persistent and should be available for all sessions
        """
        try:
            # Just remove from active sessions tracking
            # DO NOT delete the collection - it's persistent across all sessions
            with processing_lock:
                if self.session_id in active_sessions:
                    del active_sessions[self.session_id]
        except:
            pass

class SecurityError(Exception):
    """Custom exception for security-related errors"""
    pass

def get_active_sessions_info() -> Dict[str, Any]:
    """Get information about all active sessions"""
    with processing_lock:
        return {
            "total_active_sessions": len(active_sessions),
            "sessions": {
                session_id: {
                    "duration": time.time() - info["start_time"],
                    "collection_name": info["collection_name"],
                    "documents_processed": info["stats"]["documents_processed"]
                }
                for session_id, info in active_sessions.items()
            }
        }